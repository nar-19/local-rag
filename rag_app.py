import streamlit as st
import os
from io import BytesIO
from datetime import datetime
# from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from langchain_google_genai import ChatGoogleGenerativeAI # GoogleGenerativeAIEmbeddings
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, TextLoader, DirectoryLoader
from langchain_community.vectorstores import Chroma
from langchain_classic.chains import RetrievalQA
from langchain_classic import hub # For loading the RAG prompt (optional, not strictly used in this version)


# Initialize the API key from secrets file
GEMINI_API_KEY = st.secrets["API_KEY"]
# # Initialize genai client
# client = genai.Client(api_key = st.secrets["API_KEY"])
# Initialize Vector DB
VECTOR_DB_DIR = "chroma_db" # Directory to persist the Chroma DB
COLLECTION_NAME = "langchain"  # Add this constant at the top

# --- Streamlit UI Setup ---
st.set_page_config(page_title="HR Knowledge Base", layout="wide")
st.title("📚 HR Knowledge Base with RAG & LLM")
st.markdown("---")


# --- Chat bubble styling ---
st.markdown(
    """
    <style>
    .chat-row {
        display: flex;
        width: 100%;
        margin-bottom: 8px;
    }
    .chat-row.question-row {
        justify-content: flex-end;
    }
    .chat-row.answer-row {
        justify-content: flex-start;
    }
    .chat-bubble {
        max-width: 70%;
        padding: 10px 16px;
        border-radius: 16px;
        line-height: 1.4;
        word-wrap: break-word;
        box-shadow: 0 1px 2px rgba(0,0,0,0.08);
    }
    .question-bubble {
        background-color: #D6EAF8; /* light pastel blue */
        color: #1a1a1a;
        border-bottom-right-radius: 4px;
    }
    .answer-bubble {
        background-color: #FFFFC5; /* light pastel cream */
        color: #1a1a1a;
        border-bottom-left-radius: 4px;
    }
    .chat-label {
        font-size: 0.75em;
        font-weight: 600;
        opacity: 0.6;
        margin-bottom: 2px;
    }
    .sources-line {
        font-size: 0.8em;
        opacity: 0.7;
        margin-top: 6px;
    }
    </style>
    """,
    unsafe_allow_html=True,
)


# Sidebar for API Key input and instructions
with st.sidebar:
    st.header("Configuration")
    st.markdown("1.  **Set up `GOOGLE_API_KEY` in `.env` file.**")
    st.markdown("2.  **Specify your document directory.**")
    st.markdown("3.  **App will attempt to load existing knowledge base on startup.**") # Updated instruction
    st.markdown("4.  **Click 'Process Documents' to rebuild or process a new directory.**") # Updated instruction
    st.markdown("5.  **Enter your query in the chat bar.**")

    if not GEMINI_API_KEY:
        st.warning("`GOOGLE_API_KEY` not found. Please set it in a `.env` file or directly in the script.")
        st.stop() # Stop execution if API key is missing

    st.subheader("Document Processing")
    # Ensure doc_dir has a default value for initial load
    doc_dir = st.text_input("Path to documents directory:", value="./docs")
    process_button = st.button("Process Documents (Rebuild/Process New Dir)") # Updated button label

# --- Initialize Session State ---
if 'vector_store' not in st.session_state:
    st.session_state['vector_store'] = None
if 'rag_chain' not in st.session_state:
    st.session_state['rag_chain'] = None
if 'knowledge_tree_structure' not in st.session_state:
    st.session_state['knowledge_tree_structure'] = "No documents processed yet."
# New: Flag to ensure initial load/process attempt only happens once per session
if 'initial_load_attempted' not in st.session_state:
    st.session_state['initial_load_attempted'] = False
# New: Persistent chat history of question/answer pairs for this session
if 'chat_history' not in st.session_state:
    st.session_state['chat_history'] = []

# --- Functions ---

def load_documents(directory):
    """Loads documents from the specified directory."""
    documents = []
    
    try:
        st.info(f"Loading documents from {directory}...")
        
        loaders = {
            ".pdf": PyPDFLoader,
            ".txt": TextLoader,
            ".md": TextLoader # Markdown files can be loaded as text
        }
        
        # Walk through the directory to find supported files
        for root, _, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                file_extension = os.path.splitext(file_path)[1].lower()
                
                if file_extension in loaders:
                    try:
                        loader = loaders[file_extension](file_path)
                        documents.extend(loader.load())
                        st.text(f"Loaded: {os.path.basename(file_path)}")
                    except Exception as e:
                        st.warning(f"Could not load {file_path}: {e}")
                else:
                    st.info(f"Skipping unsupported file: {file_path}")
        
        return documents
    except Exception as e:
        st.error(f"Error loading documents from {directory}: {e}")
        return []

def process_documents_and_create_vector_store(doc_directory):
    """
    Attempts to load an existing Chroma DB. If not found or fails,
    it loads documents, splits them, creates embeddings, and stores them in Chroma DB.
    """
    # 1. Attempt to load existing vector store
    if os.path.exists(VECTOR_DB_DIR) and os.listdir(VECTOR_DB_DIR):
        try:
            # embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001", google_api_key=GEMINI_API_KEY)

            embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",
                model_kwargs={"device": "cpu"},   # use "cuda" if you have a GPU
                encode_kwargs={"normalize_embeddings": True}
            )
            
            st.session_state['vector_store'] = Chroma(persist_directory=VECTOR_DB_DIR, embedding_function=embeddings, collection_name=COLLECTION_NAME)
            
            # Verify if it actually contains data
            if st.session_state['vector_store']._collection.count() > 0: # Access internal count for verification
                st.success(f"Loaded existing Chroma DB from '{VECTOR_DB_DIR}' with {st.session_state['vector_store']._collection.count()} chunks.")
                st.session_state['knowledge_tree_structure'] = (
                    f"Loaded existing knowledge base from '{VECTOR_DB_DIR}'. "
                    f"Contains {st.session_state['vector_store']._collection.count()} chunks."
                )
                return True
            else:
                st.warning(f"Existing Chroma DB in '{VECTOR_DB_DIR}' found but appears empty. Rebuilding...")
                # Proceed to rebuild if empty
        except Exception as e:
            st.error(f"Error loading existing Chroma DB from '{VECTOR_DB_DIR}': {e}. Rebuilding...")
            # If loading fails, proceed to rebuild
    
    # 2. If no existing DB or loading failed, proceed to process new documents
    documents = load_documents(doc_directory)
    if not documents:
        st.warning("No supported documents found or loaded from the directory to build a new knowledge base.")
        return False

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    splits = text_splitter.split_documents(documents)

    if not splits:
        st.warning("No text chunks generated from documents.")
        return False

    st.info(f"Creating embeddings for {len(splits)} chunks. This might take a moment...")
    embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001", google_api_key=GEMINI_API_KEY)
    
    try:
        st.session_state['vector_store'] = Chroma.from_documents(
            documents=splits,
            embedding=embeddings,
            persist_directory=VECTOR_DB_DIR,
            collection_name=COLLECTION_NAME
        )
        # st.session_state['vector_store'].persist() # Ensures data is written to disk
        st.success(f"Knowledge base created with {len(splits)} chunks and saved to '{VECTOR_DB_DIR}'.")

        # Update knowledge tree structure display
        file_counts = {}
        for doc in documents:
            source = doc.metadata.get('source', 'Unknown')
            file_counts[source] = file_counts.get(source, 0) + 1
        
        structure_display = "## Processed Documents:\n"
        for source, count in file_counts.items():
            structure_display += f"- `{os.path.basename(source)}`\n"
        structure_display += f"\nTotal original documents loaded: {len(file_counts)}\n"
        structure_display += f"Total chunks created: {len(splits)}\n"
        st.session_state['knowledge_tree_structure'] = structure_display
        return True
    except Exception as e:
        st.error(f"Error creating vector store: {e}")
        return False

def setup_rag_chain():
    """Sets up the Retrieval Augmented Generation (RAG) chain."""
    if st.session_state['vector_store'] is None:
        # This case should ideally not be hit if initial loading/processing was successful
        st.error("Vector store not initialized. Cannot set up RAG chain.")
        return
    
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=GEMINI_API_KEY, temperature=0.2)
    # Retrieve top 5 relevant chunks by default
    retriever = st.session_state['vector_store'].as_retriever(search_kwargs={"k": 5}) 
    
    st.session_state['rag_chain'] = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff", # "stuff" means combining all retrieved documents into one prompt
        retriever=retriever,
        return_source_documents=True # Return the source documents used for the answer
    )
    st.info("RAG chain initialized and ready for queries!")


def render_question(question_text):
    """Renders a user question as a right-aligned pastel blue chat bubble."""
    st.markdown(
        f"""
        <div class="chat-row question-row">
            <div class="chat-bubble question-bubble">
                <div class="chat-label">You</div>
                {question_text}
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_answer(answer_text, sources=None):
    """Renders a bot answer as a left-aligned pastel cream chat bubble."""
    sources_html = ""
    if sources:
        sources_list = ", ".join(sources)
        sources_html = f'<div class="sources-line">Sources: {sources_list}</div>'

    st.markdown(
        f"""
        <div class="chat-row answer-row">
            <div class="chat-bubble answer-bubble">
                <div class="chat-label">Bot</div>
                {answer_text}
                {sources_html}
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def build_qa_presentation(chat_history):
    """
    Builds a PowerPoint presentation summarizing all Q&A pairs in chat_history.
    Returns the presentation as an in-memory BytesIO buffer (no disk writes).
    """
    # --- Color palette (Teal Trust, matches a knowledge-base/FAQ theme) ---
    PRIMARY = RGBColor(0x02, 0x80, 0x90)      # deep teal - title slide bg
    SECONDARY = RGBColor(0x00, 0xA8, 0x96)    # seafoam - accents
    ACCENT = RGBColor(0x02, 0xC3, 0x9A)       # mint - highlights
    DARK_TEXT = RGBColor(0x21, 0x2B, 0x2E)
    LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED_TEXT = RGBColor(0x5A, 0x6B, 0x6E)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # fully blank layout for full design control

    # --- Slide 1: Title slide (dark background) ---
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    bg.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Knowledge Base Q&A Summary"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BG
    p.font.name = "Calibri"

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Generated on {datetime.now().strftime('%B %d, %Y at %H:%M')} · {len(chat_history)} question(s)"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xE0, 0xF2, 0xF1)
    p.font.name = "Calibri"

    # --- One slide per Q&A pair ---
    for i, entry in enumerate(chat_history, start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_BG

        # Slide index badge (small circle, top-left) - visual motif carried across slides
        badge = slide.shapes.add_shape(9, Inches(0.5), Inches(0.5), Inches(0.6), Inches(0.6))  # oval
        badge.fill.solid()
        badge.fill.fore_color.rgb = SECONDARY
        badge.line.fill.background()
        badge.shadow.inherit = False
        badge_tf = badge.text_frame
        badge_tf.word_wrap = False
        bp = badge_tf.paragraphs[0]
        bp.text = str(i)
        bp.font.size = Pt(18)
        bp.font.bold = True
        bp.font.color.rgb = LIGHT_BG
        bp.alignment = PP_ALIGN.CENTER

        # Question header
        q_box = slide.shapes.add_textbox(Inches(1.35), Inches(0.45), Inches(11.0), Inches(1.0))
        tf = q_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = entry["question"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"

        # Divider line replaced with a soft background card for the answer (no accent stripes)
        card = slide.shapes.add_shape(1, Inches(0.7), Inches(1.7), Inches(11.9), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF2, 0xFB, 0xFA)  # very light mint tint
        card.line.color.rgb = RGBColor(0xDD, 0xF0, 0xEE)
        card.line.width = Pt(1)
        card.shadow.inherit = False

        # Answer label
        label_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(3), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = "ANSWER"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        p.font.name = "Calibri"

        # Answer body (truncate very long answers so text doesn't overflow the card)
        answer_text = entry["answer"] or ""
        max_chars = 900
        if len(answer_text) > max_chars:
            answer_text = answer_text[:max_chars].rsplit(" ", 1)[0] + "…"

        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(11.3), Inches(3.6))
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = answer_text
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"

        # Sources footer
        if entry.get("sources"):
            src_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5))
            tf = src_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Sources: " + ", ".join(entry["sources"])
            p.font.size = Pt(11)
            p.font.italic = True
            p.font.color.rgb = MUTED_TEXT
            p.font.name = "Calibri"

    # --- Closing remark slide ---
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    bg.shadow.inherit = False

    note_box = slide.shapes.add_textbox(Inches(1.3), Inches(2.6), Inches(10.7), Inches(2.3))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tip"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = (
        "These recurring questions and answers can be repurposed as the "
        "foundation for a self-service FAQ page — export them here, then "
        "curate and publish the most common ones for your users."
    )
    p2.font.size = Pt(16)
    p2.font.color.rgb = LIGHT_BG
    p2.font.name = "Calibri"
    p2.space_before = Pt(12)

    # --- Save to in-memory buffer (no disk writes) ---
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# --- Main App Logic ---

# --- Automatic Initial Load Logic (runs only once per Streamlit session) ---
if not st.session_state['initial_load_attempted']:
    st.session_state['initial_load_attempted'] = True # Mark as attempted
    
    with st.spinner("Attempting to load existing knowledge base from disk or process documents for the first time..."):
        # `doc_dir` will already have the default value from `st.text_input` if not changed
        if process_documents_and_create_vector_store(doc_dir):
            setup_rag_chain()
        else:
            # This message appears if no existing DB was loaded AND initial processing failed/found no docs
            st.warning("Knowledge base not loaded or created. Please ensure your document directory is correct and click 'Process Documents' to build it.")

# Process documents button handler (for explicit rebuild or processing a different directory)
if process_button:
    if doc_dir:
        with st.spinner("Processing documents and building knowledge base..."):
            # This call will re-run the processing logic, potentially rebuilding the DB
            if process_documents_and_create_vector_store(doc_dir):
                setup_rag_chain()
    else:
        st.warning("Please enter a directory path.")

# Display Knowledge Tree structure (or processing status)
st.subheader("Knowledge Tree Structure")
with st.expander("Click to view processed documents"):
    st.markdown(st.session_state['knowledge_tree_structure'])

    documents = load_documents(doc_dir)
    loaded_document_names = sorted(set(os.path.basename(doc.metadata.get('source', 'Unknown')) for doc in documents))
    for doc_name in loaded_document_names:
        st.markdown(f"- `{doc_name}`")


# Query bar
st.subheader("Ask a Question")
if st.session_state['rag_chain']:
    query = st.chat_input("Enter your query here... e.g. What is the leave policy during probation? \n" +
                          "Does the medical coverage covers dependents too? etc."
                          )

    if query:
        # Add the new question to the persistent chat history immediately,
        # with the answer left pending until it's generated below.
        st.session_state['chat_history'].append({
            "question": query,
            "answer": None,
            "sources": None,
        })

    # Render the full chat history (all previous Q&A pairs stay visible).
    for entry in st.session_state['chat_history']:
        # Question is always shown first, before its answer.
        render_question(entry["question"])

        if entry["answer"] is not None:
            # Already-answered turns: just display them.
            render_answer(entry["answer"], entry["sources"])
        else:
            # This is the newly submitted question - generate its answer now.
            with st.spinner("Searching and generating answer..."):
                try:
                    response = st.session_state['rag_chain'].invoke({"query": entry["question"]})
                    entry["answer"] = response['result']

                    sources = []
                    if response.get('source_documents'):
                        for doc in response['source_documents']:
                            sources.append(os.path.basename(doc.metadata.get('source', 'Unknown')))
                    entry["sources"] = sources

                    render_answer(entry["answer"], entry["sources"])
                except Exception as e:
                    entry["answer"] = f"An error occurred during query processing: {e}"
                    entry["sources"] = []
                    render_answer(entry["answer"])
                    st.info("Please ensure your API key is correct and the knowledge base is properly initialized.")
else:
    st.info("Knowledge base not ready. Please check configuration and process documents.")

# --- Sidebar: Export Q&A history as PowerPoint ---
# Placed at the end of the script (still renders inside the sidebar) so that
# a question answered earlier in this same run is already included in the export.
with st.sidebar:
    st.markdown("---")
    st.subheader("Export")
    st.caption(
        "Turn this session's Q&A into a PowerPoint deck. "
        "Handy as a starting point for building an FAQ page."
    )
    # Only offer the export once there's at least one answered Q&A pair
    answered_history = [
        e for e in st.session_state.get('chat_history', []) if e.get("answer")
    ]
    if answered_history:
        pptx_buffer = build_qa_presentation(answered_history)
        st.download_button(
            label="📊 Download Q&A as PowerPoint",
            data=pptx_buffer,
            file_name=f"qa_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    else:
        st.caption("Ask at least one question to enable the export.")
        
